from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from config import FIGURES_DIR, METRICS_DIR, REPORTS_DIR, ensure_directories


SLIDE_W = 13.333
SLIDE_H = 7.5
FONT_NAME = "Calibri"

NAVY = RGBColor(15, 23, 42)
BLUE = RGBColor(37, 99, 235)
SKY = RGBColor(14, 165, 233)
ORANGE = RGBColor(249, 115, 22)
SLATE = RGBColor(71, 85, 105)
TEXT = RGBColor(30, 41, 59)
WHITE = RGBColor(255, 255, 255)
BG = RGBColor(245, 247, 250)
BORDER = RGBColor(223, 229, 238)
LIGHT_BLUE = RGBColor(232, 240, 254)
LIGHT_ORANGE = RGBColor(255, 237, 213)
LIGHT_GREEN = RGBColor(220, 252, 231)


def set_run_style(run, size: int, color: RGBColor, bold: bool = False, font_name: str = FONT_NAME) -> None:
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_top_band(slide, title: str, eyebrow: str | None = None, slide_no: int | None = None) -> None:
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(SLIDE_W), Inches(0.82))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(0.14), Inches(0.12), Inches(0.54))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ORANGE
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.82), Inches(0.16), Inches(8.8), Inches(0.28))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_run_style(r, 24, WHITE, bold=True)

    if eyebrow:
        badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(10.25), Inches(0.18), Inches(2.0), Inches(0.32))
        badge.fill.solid()
        badge.fill.fore_color.rgb = BLUE
        badge.line.fill.background()
        badge.text_frame.clear()
        p2 = badge.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = eyebrow
        set_run_style(r2, 10, WHITE, bold=True)

    if slide_no is not None:
        footer = slide.shapes.add_textbox(Inches(12.35), Inches(7.0), Inches(0.45), Inches(0.22))
        p3 = footer.text_frame.paragraphs[0]
        p3.alignment = PP_ALIGN.RIGHT
        r3 = p3.add_run()
        r3.text = str(slide_no)
        set_run_style(r3, 10, SLATE, bold=True)


def add_card(slide, x: float, y: float, w: float, h: float, fill_color: RGBColor = WHITE, line_color: RGBColor = BORDER):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    return shape


def add_text_block(slide, x: float, y: float, w: float, h: float, lines: list[tuple[str, int, RGBColor, bool]], align: PP_ALIGN = PP_ALIGN.LEFT) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    for index, (text, size, color, bold) in enumerate(lines):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = text
        set_run_style(r, size, color, bold=bold)


def add_bullet_card(slide, x: float, y: float, w: float, h: float, title: str, bullets: list[str], fill_color: RGBColor = WHITE) -> None:
    add_card(slide, x, y, w, h, fill_color=fill_color)
    add_text_block(slide, x + 0.25, y + 0.18, w - 0.5, 0.4, [(title, 18, NAVY, True)])
    box = slide.shapes.add_textbox(Inches(x + 0.28), Inches(y + 0.72), Inches(w - 0.56), Inches(h - 0.9))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.space_after = Pt(9)
        paragraph.font.name = FONT_NAME
        paragraph.font.size = Pt(17)
        paragraph.font.color.rgb = TEXT


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)

    hero = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(SLIDE_W), Inches(1.45))
    hero.fill.solid()
    hero.fill.fore_color.rgb = NAVY
    hero.line.fill.background()

    add_card(slide, 0.8, 1.5, 9.5, 3.4)
    add_text_block(
        slide,
        1.1,
        1.9,
        8.9,
        2.1,
        [
            ("Finansal Varlıklarda Yön Tahmini ve Sinyal Üretimi", 28, TEXT, True),
            ("Güncel Proje Sunumu / Vize Sunumu", 21, BLUE, True),
            ("Günlük verilerle çalışan, model karşılaştırmalı ve dashboard destekli analiz sistemi", 16, SLATE, False),
            ("Amaç: Bir sonraki dönemde yönü öngörmek ve bunu ölçülebilir bir sinyale çevirmek", 15, SLATE, False),
        ],
    )

    stat = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(10.55), Inches(1.95), Inches(2.0), Inches(2.0))
    stat.fill.solid()
    stat.fill.fore_color.rgb = LIGHT_ORANGE
    stat.line.fill.background()
    add_text_block(
        slide,
        10.8,
        2.25,
        1.5,
        1.4,
        [
            ("5", 30, ORANGE, True),
            ("Varlık", 16, TEXT, True),
            ("4 model", 12, SLATE, False),
            ("Web arayüzü", 12, SLATE, False),
        ],
        align=PP_ALIGN.CENTER,
    )


def add_table_slide(prs: Presentation, title: str, dataframe: pd.DataFrame, slide_no: int, eyebrow: str | None = None, footnote: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_top_band(slide, title, eyebrow=eyebrow, slide_no=slide_no)

    rows, cols = dataframe.shape[0] + 1, dataframe.shape[1]
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.6), Inches(1.25), Inches(12.1), Inches(5.3))
    table = table_shape.table

    for col_index, column in enumerate(dataframe.columns):
        cell = table.cell(0, col_index)
        cell.text = str(column)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                set_run_style(r, 11, WHITE, bold=True)

    for row_index, (_, row) in enumerate(dataframe.iterrows(), start=1):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_index % 2 == 1 else RGBColor(237, 242, 247)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    set_run_style(r, 10, TEXT, bold=False)

    if footnote:
        add_text_block(slide, 0.75, 6.7, 11.5, 0.3, [(footnote, 11, SLATE, False)], align=PP_ALIGN.LEFT)


def add_chart_slide(prs: Presentation, title: str, image_path: Path, bullets: list[str], slide_no: int, eyebrow: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_top_band(slide, title, eyebrow=eyebrow, slide_no=slide_no)
    slide.shapes.add_picture(str(image_path), Inches(0.7), Inches(1.3), width=Inches(7.1), height=Inches(4.9))
    add_bullet_card(slide, 8.15, 1.35, 4.45, 4.8, "Kısa Yorum", bullets, fill_color=WHITE)


def add_two_column_slide(
    prs: Presentation,
    title: str,
    left_title: str,
    left_bullets: list[str],
    right_title: str,
    right_bullets: list[str],
    slide_no: int,
    eyebrow: str | None = None,
    footer_text: str | None = None,
    footer_link: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_top_band(slide, title, eyebrow=eyebrow, slide_no=slide_no)
    add_bullet_card(slide, 0.7, 1.35, 5.9, 5.0, left_title, left_bullets, fill_color=WHITE)
    add_bullet_card(slide, 6.8, 1.35, 5.85, 5.0, right_title, right_bullets, fill_color=LIGHT_BLUE)

    if footer_text:
        footer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(6.55), Inches(11.35), Inches(0.45))
        footer.fill.solid()
        footer.fill.fore_color.rgb = RGBColor(233, 239, 247)
        footer.line.fill.background()
        footer.text_frame.clear()
        p = footer.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = footer_text
        set_run_style(r, 12, BLUE, bold=True)
        if footer_link:
            r.hyperlink.address = footer_link


def load_frames() -> tuple[pd.DataFrame, pd.DataFrame]:
    recommendations = pd.read_csv(METRICS_DIR / "final_presentation_recommendations.csv")
    recent_summary = pd.read_csv(METRICS_DIR / "final_recent_signal_summary.csv")
    return recommendations, recent_summary


def build_holdout_table(recommendations: pd.DataFrame) -> pd.DataFrame:
    frame = recommendations[["symbol", "best_holdout_model", "holdout_roc_auc"]].copy()
    frame.columns = ["Varlık", "En İyi Model", "ROC-AUC"]
    frame["ROC-AUC"] = frame["ROC-AUC"].map(lambda v: f"{v:.4f}")
    return frame.sort_values("Varlık").reset_index(drop=True)


def build_recommendation_table(recommendations: pd.DataFrame) -> pd.DataFrame:
    frame = recommendations[["symbol", "recommended_model", "recommended_total_return", "recommended_sharpe"]].copy()
    frame.columns = ["Varlık", "Önerilen Model", "Toplam Getiri", "Sharpe"]
    frame["Toplam Getiri"] = frame["Toplam Getiri"].map(lambda v: f"{v:.4f}")
    frame["Sharpe"] = frame["Sharpe"].map(lambda v: f"{v:.4f}")
    return frame.sort_values("Varlık").reset_index(drop=True)


def build_recent_signal_table(recent_summary: pd.DataFrame) -> pd.DataFrame:
    frame = recent_summary[["symbol", "model", "observation_count", "active_trade_days", "up_signals", "down_signals", "no_trade_signals", "active_hit_rate", "latest_signal", "latest_probability"]].copy()
    frame.columns = ["Varlık", "Model", "Gözlem", "Aktif Gün", "Yukarı", "Aşağı", "İşlem Yok", "Aktif İşlem Başarısı", "Son Sinyal", "Son Olasılık"]
    frame["Aktif İşlem Başarısı"] = frame["Aktif İşlem Başarısı"].map(lambda v: "-" if pd.isna(v) else f"{v:.2%}")
    frame["Son Olasılık"] = frame["Son Olasılık"].map(lambda v: f"{v:.4f}")
    return frame.sort_values("Varlık").reset_index(drop=True)


def create_presentation() -> Path:
    ensure_directories()
    recommendations, recent_summary = load_frames()
    window_label = "Son 10 Takvim Günü"
    if not recent_summary.empty:
        window_label = f"Son 10 Takvim Günü ({recent_summary.iloc[0]['window_start']} - {recent_summary.iloc[0]['window_end']})"
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    add_title_slide(prs)
    add_two_column_slide(
        prs,
        "Problem ve Hedef",
        "Neyi Çözmeye Çalışıyor?",
        [
            "Bu çalışma, fiyatın tam değerini tahmin etmeye çalışmıyor.",
            "Amaç, bir sonraki dönemde yönün yukarı mı aşağı mı olacağını kestirmek.",
            "Böylece model çıktısı daha sade bir finansal sinyale dönüşebiliyor.",
        ],
        "Neden Önemli?",
        [
            "Finansal piyasalarda her zaman tam fiyat tahmini gerekmez.",
            "Bazen doğru yön bilgisi bile karar kalitesini artırabilir.",
            "Bu proje de tam olarak bu noktaya odaklanıyor.",
        ],
        slide_no=2,
        eyebrow="Genel Bakış",
    )
    add_two_column_slide(
        prs,
        "Sistem Ne Yapıyor?",
        "Kapsam",
        [
            "5 varlık: BTC/USD, ETH/USD, AAPL, NVDA, XAU/USD",
            "Günlük veri kullanılıyor",
            "Birden fazla model aynı veri hattısında karşılaştırılıyor",
            "Çıktı sadece tahmin değil, sinyal ve backtest sonucu da üretiyor",
        ],
        "Üretilen Çıktılar",
        [
            "Yukarı / Aşağı / İşlem Yok sinyali",
            "Holdout sonuçları",
            "Backtest sonuçları",
            "Rolling retraining karşılaştırması",
            "Web tabanlı izleme arayüzü",
        ],
        slide_no=3,
        eyebrow="Sistem Özeti",
    )
    add_two_column_slide(
        prs,
        "Nasıl Çalışıyor?",
        "Veri Hattısı",
        [
            "API üzerinden veri çekilir",
            "Teknik göstergeler ve bağlamsal özellikler üretilir",
            "Veri sızıntısını engellemek için zaman bazlı bölme yapılır",
            "Farklı modeller aynı yapı üzerinde eğitilir",
        ],
        "Karar Katmanı",
        [
            "Model olasılık üretir",
            "Bu olasılık doğrudan işleme dönüşmez",
            "Güven düşükse İşlem Yok seçeneği devreye girer",
            "Böylece sistem daha gerçekçi davranır",
        ],
        slide_no=4,
        eyebrow="Yöntem",
    )
    add_two_column_slide(
        prs,
        "Kullanılan Modeller",
        "Klasik Başlangıç Modelleri",
        [
            "Logistic Regression: en temel karşılaştırma noktasıdır.",
            "MLP: klasik yöntem ile derin öğrenme arasında köprü kuran daha esnek bir yapıdır.",
            "Bu iki model, sistemin basit ve anlaşılır tarafını temsil eder.",
        ],
        "Derin Öğrenme Modelleri",
        [
            "LSTM: zaman sıralı verilerde geçmiş bilgiyi daha iyi taşıyabilmek için kullanılır.",
            "GRU: LSTM'e benzer ama daha sade bir sıralı modeldir.",
            "Amaç, aynı veri üzerinde farklı model ailelerini karşılaştırmaktır.",
        ],
        slide_no=5,
        eyebrow="Model Katmanı",
    )
    add_two_column_slide(
        prs,
        "Metrikler Ne Anlatıyor?",
        "Tahmin Başarısı",
        [
            "Accuracy: genel doğru tahmin oranı",
            "F1: yukarı/aşağı dengesini birlikte ölçen daha dengeli metrik",
            "ROC-AUC: modelin yönleri birbirinden ayırma gücü",
        ],
        "Finansal Yorum",
        [
            "Toplam Getiri: stratejinin test dönemindeki bileşik sonucu",
            "Sharpe: getiriyi risk ile birlikte yorumlayan ölçü",
            "Coverage: modelin gerçekten işlem önerdiği gün oranı",
        ],
        slide_no=6,
        eyebrow="Metrikler",
    )
    add_table_slide(
        prs,
        "Model Performans Özeti",
        build_holdout_table(recommendations),
        slide_no=7,
        eyebrow="Holdout",
        footnote="Buradaki tablo her varlıkta test döneminde en yüksek ROC-AUC veren modeli özetler.",
    )
    add_chart_slide(
        prs,
        "Görsel Sonuç: Holdout Başarısı",
        FIGURES_DIR / "final_holdout_best_roc_auc.png",
        [
            "ETH/USD ve NVDA tarafında göreli olarak daha güçlü ayrıştırma görüyoruz.",
            "Değerler hâlâ orta seviyede; yani sistem temkinli okunmalı.",
            "Bu yüzden proje sadece skor raporlamak yerine backtest ve sinyal katmanı da içeriyor.",
        ],
        slide_no=8,
        eyebrow="Teknik Sonuç",
    )
    add_table_slide(
        prs,
        "Önerilen Modeller",
        build_recommendation_table(recommendations),
        slide_no=9,
        eyebrow="Karar Katmanı",
        footnote="Toplam Getiri, test döneminde önerilen model ve no-trade bandı ile oluşan stratejinin bileşik getirisidir. Örnek: 1.10 değeri yaklaşık %110 toplam artış anlamına gelir. Sharpe ise getiriyi oynaklığa göre düzelten risk-ayarlı performans göstergesidir; daha yüksek olması daha dengeli sonuç anlamına gelir.",
    )
    add_chart_slide(
        prs,
        "Görsel Sonuç: Önerilen Modellerin Getirisi",
        FIGURES_DIR / "final_recommended_total_return.png",
        [
            "XAU/USD ve NVDA tarafı finansal olarak daha güçlü sinyal üretti.",
            "BTC/USD daha zor bir piyasa; bu yüzden karar katmanı çok daha kritik.",
            "Buradaki sonuçlar model çıktısını doğrudan sinyale çevirdiğimizde ne olduğunu gösteriyor.",
        ],
        slide_no=10,
        eyebrow="Finansal Sonuç",
    )
    add_table_slide(
        prs,
        window_label,
        build_recent_signal_table(recent_summary),
        slide_no=11,
        eyebrow="Örnek Çıktı",
        footnote="Bu tabloda tarih araligi baslikta verilir. Aktif Islem Basarisi yalnizca islem acilan gunler uzerinden hesaplanir; ornegin 9 gun Islem Yok ve 1 dogru sinyal varsa bu oran %100 gorunebilir. XAU/USD tarafinda veri saglayici hafta sonu tarihli gunluk barlar dondurebildigi icin 10 gozlem gorulmektedir.",
    )
    add_chart_slide(
        prs,
        "Son 10 Gün Sinyal Dağılımı",
        FIGURES_DIR / "final_recent_signal_distribution.png",
        [
            "Bazı varlıklarda model sık sinyal üretirken, bazılarında daha seçici davranıyor.",
            "İşlem Yok alanı, düşük güvenli günleri filtrelemek için önemli.",
            "Bu yaklaşım sistemi daha kontrollü hale getiriyor.",
        ],
        slide_no=12,
        eyebrow="Örnekler",
    )
    add_chart_slide(
        prs,
        "Statik ve Zaman İçinde Güncellenen Model Karşılaştırması",
        FIGURES_DIR / "final_rolling_vs_recommended.png",
        [
            "Rolling retraining, modeli zaman içinde yeniden eğiterek daha gerçekçi bir test sağlar.",
            "BTC/USD ve NVDA için bu karşılaştırmayı yaptık.",
            "Özellikle NVDA tarafında rolling yapı daha güçlü bir potansiyel gösteriyor.",
        ],
        slide_no=13,
        eyebrow="Rolling Retraining",
    )
    add_two_column_slide(
        prs,
        "Neler Öğrendik?",
        "Güçlü Taraflar",
        [
            "Sistem artık sadece model eğiten bir yapı değil, sinyal üreten bir yapıya dönüştü.",
            "Aynı anda birden fazla varlık ve model karşılaştırılabiliyor.",
            "Web arayüzü üzerinden sonuçlar daha anlaşılır hale geldi.",
        ],
        "Dikkat Edilmesi Gerekenler",
        [
            "Bu sonuçlar geçmiş veriye dayalıdır; gelecekte aynı sonucu garanti etmez.",
            "Bazı varlıklar diğerlerine göre çok daha zor davranıyor.",
            "Backtest mantığı güçlü bir araçtır ama her zaman risk yönetimi ile birlikte okunmalıdır.",
        ],
        slide_no=14,
        eyebrow="Yorum",
    )
    add_two_column_slide(
        prs,
        "Sonuç ve Sonraki Adım",
        "Bugünkü Durum",
        [
            "Bu sunum projenin bugünkü durumunu özetler; final sunum proje tamamlandığında ayrıca hazırlanacaktır.",
            "Veri hattısı, modelleme, backtest ve dashboard katmanları birlikte çalışıyor.",
            "Her varlık için önerilen model ve sinyal mantığı görülebiliyor.",
        ],
        "Sonraki Doğal Adımlar",
        [
            "Dashboard içinden daha güçlü yönetim butonları",
            "Rolling retraining'i daha fazla modele yaymak",
            "Canlıya yakın veri güncelleme katmanı",
            "Paper trading ve daha gelişmiş risk yönetimi",
        ],
        slide_no=15,
        eyebrow="Kapanış",
        footer_text="GitHub Repo: https://github.com/furkann44/finDL",
        footer_link="https://github.com/furkann44/finDL",
    )

    output_path = REPORTS_DIR / "vize_sunumu_guncel.pptx"
    try:
        prs.save(output_path)
        return output_path
    except PermissionError:
        fallback_path = REPORTS_DIR / "vize_sunumu_guncel_updated.pptx"
        prs.save(fallback_path)
        return fallback_path


def main() -> None:
    output_path = create_presentation()
    print(output_path)


if __name__ == "__main__":
    main()
