from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from config import FIGURES_DIR, METRICS_DIR, PREDICTIONS_DIR, REPORTS_DIR, ensure_directories


SLIDE_W = 13.333
SLIDE_H = 7.5
FONT_NAME = "Calibri"

NAVY = RGBColor(15, 23, 42)
BLUE = RGBColor(37, 99, 235)
SKY = RGBColor(14, 165, 233)
ORANGE = RGBColor(249, 115, 22)
SLATE = RGBColor(71, 85, 105)
TEXT = RGBColor(30, 41, 59)
WHITE = RGBColor(255, 255, 255)
BG = RGBColor(245, 247, 250)
BORDER = RGBColor(223, 229, 238)
LIGHT_BLUE = RGBColor(232, 240, 254)
LIGHT_ORANGE = RGBColor(255, 237, 213)
LIGHT_GREEN = RGBColor(220, 252, 231)


def set_run_style(run, size: int, color: RGBColor, bold: bool = False, font_name: str = FONT_NAME) -> None:
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_top_band(slide, title: str, eyebrow: str | None = None, slide_no: int | None = None) -> None:
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(SLIDE_W), Inches(0.82))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(0.14), Inches(0.12), Inches(0.54))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ORANGE
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.82), Inches(0.16), Inches(8.4), Inches(0.28))
    title_frame = title_box.text_frame
    title_frame.clear()
    paragraph = title_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = title
    set_run_style(run, 24, WHITE, bold=True)

    if eyebrow:
        badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(10.15), Inches(0.18), Inches(1.95), Inches(0.32))
        badge.fill.solid()
        badge.fill.fore_color.rgb = BLUE
        badge.line.fill.background()
        badge.text_frame.clear()
        p = badge.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = eyebrow
        set_run_style(r, 10, WHITE, bold=True)

    if slide_no is not None:
        footer = slide.shapes.add_textbox(Inches(12.35), Inches(7.0), Inches(0.45), Inches(0.22))
        p = footer.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = str(slide_no)
        set_run_style(r, 10, SLATE, bold=True)


def add_card(slide, x: float, y: float, w: float, h: float, fill_color: RGBColor = WHITE, line_color: RGBColor = BORDER):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = line_color
    return card


def add_text_block(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    lines: list[tuple[str, int, RGBColor, bool]],
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP

    for index, (text, size, color, bold) in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = align
        paragraph.space_after = Pt(4)
        run = paragraph.add_run()
        run.text = text
        set_run_style(run, size, color, bold=bold)


def add_bullet_card(slide, x: float, y: float, w: float, h: float, title: str, bullets: list[str], fill_color: RGBColor = WHITE) -> None:
    add_card(slide, x, y, w, h, fill_color=fill_color)
    add_text_block(slide, x + 0.25, y + 0.18, w - 0.5, 0.4, [(title, 18, NAVY, True)])

    box = slide.shapes.add_textbox(Inches(x + 0.28), Inches(y + 0.72), Inches(w - 0.56), Inches(h - 0.9))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.space_after = Pt(10)
        paragraph.font.name = FONT_NAME
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = TEXT


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)

    hero = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(SLIDE_W), Inches(1.45))
    hero.fill.solid()
    hero.fill.fore_color.rgb = NAVY
    hero.line.fill.background()

    circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(11.0), Inches(0.85), Inches(1.5), Inches(1.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ORANGE
    circle.line.fill.background()

    circle2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(10.25), Inches(1.55), Inches(0.85), Inches(0.85))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = SKY
    circle2.line.fill.background()

    add_card(slide, 0.8, 1.55, 9.35, 3.2)
    add_text_block(
        slide,
        1.15,
        1.95,
        8.6,
        1.8,
        [
            ("Finansal Piyasalarda Derin Öğrenme ile Yön Tahmini", 28, TEXT, True),
            ("Vize Sunumu", 22, BLUE, True),
            ("Yüksek Lisans Derin Öğrenme Dersi", 15, SLATE, False),
        ],
    )

    tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.15), Inches(4.2), Inches(5.5), Inches(0.45))
    tag.fill.solid()
    tag.fill.fore_color.rgb = LIGHT_BLUE
    tag.line.fill.background()
    tag.text_frame.clear()
    p = tag.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Kapsam: veri hattı, feature engineering, leakage önleme, baseline ve ilk LSTM sonuçları"
    set_run_style(r, 12, NAVY, bold=True)

    add_card(slide, 10.45, 2.35, 2.2, 2.0, fill_color=LIGHT_ORANGE, line_color=LIGHT_ORANGE)
    add_text_block(
        slide,
        10.75,
        2.72,
        1.6,
        1.2,
        [
            ("5", 28, ORANGE, True),
            ("Varlık", 16, TEXT, True),
            ("1 günlük veri", 12, SLATE, False),
        ],
        align=PP_ALIGN.CENTER,
    )


def add_problem_slide(prs: Presentation, slide_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_top_band(slide, "Problem Tanımı", eyebrow="Vize Kapsamı", slide_no=slide_no)

    add_bullet_card(
        slide,
        0.7,
        1.25,
        5.9,
        4.9,
        "Temel Soru",
        [
            "Amaç, fiyatın tam değerini değil bir sonraki periyottaki yönünü tahmin etmek.",
            "Hedef değişken: close(t+1) > close(t) ise 1, aksi halde 0.",
            "Problem tipi: zaman serisi tabanlı ikili sınıflandırma.",
        ],
        fill_color=WHITE,
    )
    add_bullet_card(
        slide,
        6.85,
        1.25,
        5.8,
        4.9,
        "Bu Aşamadaki Öncelik",
        [
            "En yüksek skoru zorlamak değil, güvenilir bir deney hattısı kurmak.",
            "Veri sızıntısını engellemek.",
            "Uçtan uca çalışan ilk derin öğrenme prototipini göstermek.",
        ],
        fill_color=LIGHT_BLUE,
    )


def add_scope_slide(prs: Presentation, slide_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_top_band(slide, "Veri ve Kapsam", eyebrow="Kurulum", slide_no=slide_no)

    add_bullet_card(
        slide,
        0.7,
        1.25,
        4.9,
        5.15,
        "Veri Kaynağı ve Parametreler",
        [
            "Veri kaynağı: Twelve Data API",
            "Frekans: 1day",
            "Timezone: UTC",
            "Lookback: 30",
            "İlk aşamada canlı veri kullanılmadı",
        ],
    )

    add_card(slide, 5.95, 1.25, 6.7, 5.15, fill_color=WHITE)
    add_text_block(slide, 6.25, 1.48, 3.8, 0.35, [("İncelenen Varlıklar", 18, NAVY, True)])

    tags = [
        ("BTC/USD", LIGHT_BLUE, BLUE),
        ("ETH/USD", LIGHT_BLUE, BLUE),
        ("AAPL", LIGHT_GREEN, RGBColor(22, 101, 52)),
        ("NVDA", LIGHT_GREEN, RGBColor(22, 101, 52)),
        ("XAU/USD", LIGHT_ORANGE, ORANGE),
    ]
    positions = [(6.25, 2.0), (8.2, 2.0), (10.15, 2.0), (6.25, 2.75), (8.2, 2.75)]
    for (label, fill_color, text_color), (x, y) in zip(tags, positions, strict=True):
        tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.65), Inches(0.45))
        tag.fill.solid()
        tag.fill.fore_color.rgb = fill_color
        tag.line.fill.background()
        tag.text_frame.clear()
        p = tag.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        set_run_style(r, 12, text_color, bold=True)

    add_text_block(
        slide,
        6.25,
        3.7,
        5.8,
        1.9,
        [
            ("Vize sınırı", 16, NAVY, True),
            ("Günlük veri, feature engineering, Logistic Regression baseline ve ilk LSTM sonuçları.", 18, TEXT, False),
        ],
    )


def add_pipeline_slide(prs: Presentation, slide_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_top_band(slide, "Sistem Mimarisi", eyebrow="Veri Hattısı", slide_no=slide_no)

    labels = [
        "API'den\nveri çek",
        "Feature\nengineering",
        "Target\nüretimi",
        "Zaman bazlı\nsplit",
        "Scaling\n(train only)",
        "Baseline +\nLSTM",
        "Metrikler\nve kayıt",
    ]
    x_positions = [0.55, 2.35, 4.15, 5.95, 7.75, 9.55, 11.35]

    for index, (label, x) in enumerate(zip(labels, x_positions, strict=True)):
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.6), Inches(1.35), Inches(1.1))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE if index % 2 == 0 else LIGHT_BLUE
        box.line.color.rgb = BORDER
        box.text_frame.clear()
        p = box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        set_run_style(r, 14, TEXT, bold=True)

        if index < len(labels) - 1:
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(x + 1.35),
                Inches(3.15),
                Inches(x_positions[index + 1]),
                Inches(3.15),
            )
            connector.line.color.rgb = BLUE
            connector.line.width = Pt(2)

    add_text_block(
        slide,
        0.85,
        4.55,
        11.5,
        1.3,
        [
            ("Akışın amacı, ham veriyi model girdi formatına dönüştürmek ve her aşamada veri sızıntısını kontrol altında tutmaktır.", 18, TEXT, False),
        ],
        align=PP_ALIGN.CENTER,
    )


def add_features_leakage_slide(prs: Presentation, slide_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_top_band(slide, "Özellikler ve Leakage Önleme", eyebrow="Metodoloji", slide_no=slide_no)

    add_bullet_card(
        slide,
        0.7,
        1.25,
        6.0,
        5.2,
        "Üretilen Özellikler",
        [
            "Getiri: return_1, return_5, co_return, hl_spread",
            "Volatilite: volatility_10",
            "Trend: sma_5, sma_10, ema_10",
            "Momentum: rsi_14, macd, macd_signal, macd_hist",
            "Etiket: future_return_1 ve target",
        ],
        fill_color=WHITE,
    )
    add_bullet_card(
        slide,
        6.95,
        1.25,
        5.7,
        5.2,
        "Leakage Önleme Kuralları",
        [
            "Bölme işlemi tamamen zamana göre yapıldı",
            "Shuffle kullanılmadı",
            "StandardScaler yalnızca train set üzerinde fit edildi",
            "Geleceğe ait bilgi feature içinde kullanılmadı",
            "t+1 bilgisi yalnızca hedef değişkende yer aldı",
        ],
        fill_color=LIGHT_ORANGE,
    )


def add_completed_slide(prs: Presentation, slide_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_top_band(slide, "Vize Kapsamında Tamamlananlar", eyebrow="İlerleme", slide_no=slide_no)

    checkpoints = [
        "5 varlık için günlük ham veri çekildi",
        "İşlenmiş veri setleri üretildi",
        "Teknik göstergeler ve hedef etiketi oluşturuldu",
        "Logistic Regression baseline kuruldu",
        "LSTM ile ilk derin öğrenme denemesi yapıldı",
        "Metrikler ve deney çıktıları otomatik kaydedildi",
    ]

    for index, text in enumerate(checkpoints):
        y = 1.35 + index * 0.78
        card = add_card(slide, 1.0, y, 11.2, 0.56, fill_color=WHITE)
        icon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(1.18), Inches(y + 0.11), Inches(0.28), Inches(0.28))
        icon.fill.solid()
        icon.fill.fore_color.rgb = BLUE
        icon.line.fill.background()
        add_text_block(slide, 1.58, y + 0.1, 10.0, 0.3, [(text, 18, TEXT, False)])


def build_results_table() -> pd.DataFrame:
    summary = pd.read_csv(METRICS_DIR / "model_summary.csv")
    filtered = summary[(summary["split"] == "test") & (summary["model"].isin(["baseline", "lstm"]))][
        ["symbol", "model", "roc_auc"]
    ]
    pivot = filtered.pivot(index="symbol", columns="model", values="roc_auc").reset_index()
    pivot.columns = ["Varlık", "Baseline ROC-AUC", "LSTM ROC-AUC"]
    pivot["Baseline ROC-AUC"] = pivot["Baseline ROC-AUC"].map(lambda value: f"{value:.4f}")
    pivot["LSTM ROC-AUC"] = pivot["LSTM ROC-AUC"].map(lambda value: f"{value:.4f}")
    return pivot


def direction_label(prediction: int) -> str:
    return "Yukarı" if int(prediction) == 1 else "Aşağı"


def build_signal_snapshot_table() -> pd.DataFrame:
    summary = pd.read_csv(METRICS_DIR / "model_summary.csv")
    vize_models = summary[(summary["split"] == "test") & (summary["model"].isin(["baseline", "lstm"]))][
        ["symbol", "model", "roc_auc"]
    ]
    best = vize_models.sort_values(["symbol", "roc_auc"], ascending=[True, False]).drop_duplicates("symbol")

    rows: list[dict[str, str]] = []
    for row in best.itertuples(index=False):
        file_name = f"{row.symbol.lower().replace('/', '_')}_{row.model}_test_predictions.csv"
        prediction_frame = pd.read_csv(PREDICTIONS_DIR / file_name)
        latest = prediction_frame.iloc[-1]
        rows.append(
            {
                "Varlık": row.symbol,
                "Model": str(row.model).upper(),
                "Tarih": str(latest["datetime"])[:10],
                "Olasılık": f"{float(latest['probability']):.4f}",
                "Tahmin": direction_label(int(latest["prediction"])),
            }
        )

    return pd.DataFrame(rows)


def add_results_slide(prs: Presentation, slide_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_top_band(slide, "Ara Sonuçlar", eyebrow="Baseline vs LSTM", slide_no=slide_no)

    table_frame = build_results_table()
    rows, cols = table_frame.shape[0] + 1, table_frame.shape[1]
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.7), Inches(1.35), Inches(8.05), Inches(4.9))
    table = table_shape.table

    for col_index, column in enumerate(table_frame.columns):
        cell = table.cell(0, col_index)
        cell.text = str(column)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                set_run_style(run, 12, WHITE, bold=True)

    for row_index, (_, row) in enumerate(table_frame.iterrows(), start=1):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_index % 2 == 1 else RGBColor(237, 242, 247)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    set_run_style(run, 12, TEXT, bold=False)

    add_bullet_card(
        slide,
        9.15,
        1.35,
        3.5,
        4.9,
        "Hızlı Okuma",
        [
            "ETH/USD tarafında LSTM, baseline'a göre daha güçlü ROC-AUC verdi.",
            "XAU/USD tarafında LSTM küçük bir iyileşme sağladı.",
            "BTC/USD tarafında LSTM henüz baseline'ı geçemedi.",
        ],
        fill_color=LIGHT_BLUE,
    )


def add_signal_slide(prs: Presentation, slide_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_top_band(slide, "Yukarı / Aşağı Tahmini", eyebrow="Finansal Sinyal", slide_no=slide_no)

    add_bullet_card(
        slide,
        0.7,
        1.35,
        5.25,
        4.8,
        "Model Çıktısı Nasıl Yorumlanır?",
        [
            "Model her örnek için bir yukarı yön olasılığı üretir.",
            "Olasılık 0.50 üzerindeyse Yukarı, altında ise Aşağı sinyali verilir.",
            "Bu yapı, projeyi sadece sınıflandırma olmaktan çıkarıp finansal analiz yönüne taşır.",
            "Gerçek kullanımda buna bir de İşlem Yok bölgesi eklenebilir.",
        ],
        fill_color=WHITE,
    )

    snapshot = build_signal_snapshot_table()
    rows, cols = snapshot.shape[0] + 1, snapshot.shape[1]
    table_shape = slide.shapes.add_table(rows, cols, Inches(6.25), Inches(1.55), Inches(6.1), Inches(3.7))
    table = table_shape.table

    for col_index, column in enumerate(snapshot.columns):
        cell = table.cell(0, col_index)
        cell.text = str(column)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                set_run_style(run, 11, WHITE, bold=True)

    for row_index, (_, row) in enumerate(snapshot.iterrows(), start=1):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = str(value)
            cell.fill.solid()
            if snapshot.columns[col_index] == "Tahmin":
                cell.fill.fore_color.rgb = LIGHT_GREEN if value == "Yukarı" else LIGHT_ORANGE
            else:
                cell.fill.fore_color.rgb = WHITE if row_index % 2 == 1 else RGBColor(237, 242, 247)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    set_run_style(run, 11, TEXT, bold=snapshot.columns[col_index] == "Tahmin")

    add_text_block(
        slide,
        6.3,
        5.55,
        5.95,
        0.7,
        [
            ("Not: Bu tablo canlı sinyal değil, test setindeki en güncel örnek için model çıktısını göstermektedir.", 11, SLATE, False),
        ],
    )


def add_chart_slide(prs: Presentation, slide_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_top_band(slide, "Görsel Özet", eyebrow="ROC-AUC", slide_no=slide_no)

    chart_path = FIGURES_DIR / "vize_baseline_lstm_test_roc_auc.png"
    slide.shapes.add_picture(str(chart_path), Inches(0.7), Inches(1.35), width=Inches(7.4), height=Inches(4.7))

    add_bullet_card(
        slide,
        8.4,
        1.35,
        4.2,
        4.7,
        "Gözlemler",
        [
            "ETH/USD ve XAU/USD tarafında LSTM ayrıştırma gücünü bir miktar artırdı.",
            "BTC/USD için LSTM henüz baseline'ın önüne geçemedi.",
            "Sonuçlar umut verici ancak hâlâ orta seviyede.",
        ],
        fill_color=WHITE,
    )


def add_findings_slide(prs: Presentation, slide_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_top_band(slide, "İlk Teknik Yorumlar", eyebrow="Ara Değerlendirme", slide_no=slide_no)

    add_bullet_card(
        slide,
        0.7,
        1.5,
        3.85,
        4.55,
        "Olumlu Taraflar",
        [
            "Veri hattısı artık stabil ve tekrar üretilebilir.",
            "Beş farklı varlıkta aynı akış çalışıyor.",
            "Baseline ve derin öğrenme karşılaştırması mümkün hale geldi.",
        ],
        fill_color=LIGHT_GREEN,
    )
    add_bullet_card(
        slide,
        4.75,
        1.5,
        3.85,
        4.55,
        "Temel Gözlem",
        [
            "Model davranışı piyasa tipine göre değişiyor.",
            "Kripto ve emtia tarafı, hisse senedinden farklı tepki veriyor.",
            "Bu nedenle tek bir modelle tüm piyasaları açıklamak zor.",
        ],
        fill_color=LIGHT_BLUE,
    )
    add_bullet_card(
        slide,
        8.8,
        1.5,
        3.85,
        4.55,
        "Sınırlılıklar",
        [
            "ROC-AUC değerleri genel olarak orta seviyede kaldı.",
            "Bu da ek geliştirmelerin gerekli olduğunu gösteriyor.",
            "Vizede amaç performans zirvesi değil, doğru metodolojiyi göstermektir.",
        ],
        fill_color=LIGHT_ORANGE,
    )


def add_next_steps_slide(prs: Presentation, slide_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_top_band(slide, "Finale Kalanlar", eyebrow="Sonraki Aşama", slide_no=slide_no)

    add_bullet_card(
        slide,
        0.9,
        1.45,
        11.4,
        4.9,
        "Final Aşamasında Eklenecekler",
        [
            "MLP ve GRU ile daha geniş model karşılaştırması",
            "Walk-forward validation",
            "Threshold-based labeling",
            "Basit backtest",
            "4 saatlik veri ile genişleme",
            "ROC / confusion matrix görselleri ve final raporu",
        ],
        fill_color=WHITE,
    )


def add_closing_slide(prs: Presentation, slide_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

    add_text_block(
        slide,
        1.0,
        1.5,
        9.2,
        2.0,
        [
            ("Sonuç", 28, WHITE, True),
            ("Veri toplama, feature engineering ve temel modelleme başarıyla tamamlandı.", 20, WHITE, False),
            ("Proje, final aşamasına taşınabilecek sağlam bir temel kazandı.", 20, WHITE, False),
        ],
    )

    ribbon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.7), Inches(3.2), Inches(0.6))
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = ORANGE
    ribbon.line.fill.background()
    ribbon.text_frame.clear()
    p = ribbon.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Teşekkürler"
    set_run_style(r, 18, WHITE, bold=True)

    footer = slide.shapes.add_textbox(Inches(10.7), Inches(6.55), Inches(1.6), Inches(0.25))
    p2 = footer.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = str(slide_no)
    set_run_style(r2, 10, RGBColor(226, 232, 240), bold=True)


def create_presentation() -> Path:
    ensure_directories()
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    add_title_slide(prs)
    add_problem_slide(prs, 2)
    add_scope_slide(prs, 3)
    add_pipeline_slide(prs, 4)
    add_features_leakage_slide(prs, 5)
    add_completed_slide(prs, 6)
    add_results_slide(prs, 7)
    add_signal_slide(prs, 8)
    add_chart_slide(prs, 9)
    add_findings_slide(prs, 10)
    add_next_steps_slide(prs, 11)
    add_closing_slide(prs, 12)

    output_path = REPORTS_DIR / "vize_presentation.pptx"
    try:
        prs.save(output_path)
        return output_path
    except PermissionError:
        fallback_path = REPORTS_DIR / "vize_presentation_updated.pptx"
        prs.save(fallback_path)
        return fallback_path


def main() -> None:
    output_path = create_presentation()
    print(output_path)


if __name__ == "__main__":
    main()
